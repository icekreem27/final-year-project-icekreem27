import os
from pptx import Presentation

def convert_pptx_folder_to_txt(source_folder_path, destination_folder_path):
    # Check if the destination folder exists, if not, create it
    if not os.path.exists(destination_folder_path):
        os.makedirs(destination_folder_path)
    
    # Iterate through all files in the folder
    for filename in os.listdir(source_folder_path):
        if filename.endswith(".pptx"):
            # Construct full file path for the pptx
            pptx_path = os.path.join(source_folder_path, filename)
            # Open the pptx file
            prs = Presentation(pptx_path)
            text = ""
            # Extract text from each slide
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            # Construct text file name and path in the destination folder
            txt_filename = os.path.splitext(filename)[0] + os.path.splitext(filename)[1] + ".txt"
            txt_path = os.path.join(destination_folder_path, txt_filename)
            
            # Save the extracted text to a .txt file
            with open(txt_path, "w", encoding="utf-8") as text_file:
                text_file.write(text)

# Set folder paths
source_folder_path = 'Data Mining files/scraped files'
destination_folder_path = 'Data Mining files/pptx_txt_files'
convert_pptx_folder_to_txt(source_folder_path, destination_folder_path)
